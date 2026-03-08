from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Notebook Tabs inspired palette
COLOR_BG = RGBColor(0x27, 0x28, 0x22)
COLOR_PAPER = RGBColor(0xF6, 0xF3, 0xEA)
COLOR_INK = RGBColor(0x1B, 0x1D, 0x1F)
COLOR_SOFT = RGBColor(0x4E, 0x52, 0x57)
COLOR_ACCENT = RGBColor(0x2F, 0x6E, 0x9B)
COLOR_ACCENT_2 = RGBColor(0x2F, 0x8F, 0x6E)
COLOR_TAB_1 = RGBColor(0x95, 0xD4, 0xB7)
COLOR_TAB_2 = RGBColor(0xC5, 0xB4, 0xE8)
COLOR_TAB_3 = RGBColor(0xF2, 0xB2, 0xC2)
COLOR_TAB_4 = RGBColor(0xA8, 0xD4, 0xEA)
COLOR_TAB_5 = RGBColor(0xFF, 0xE2, 0xA2)


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_sheet(slide):
    sheet = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(0.45),
        Inches(12.35),
        Inches(6.6),
    )
    sheet.fill.solid()
    sheet.fill.fore_color.rgb = COLOR_PAPER
    sheet.line.color.rgb = RGBColor(0xE3, 0xDE, 0xD2)
    sheet.line.width = Pt(1)

    # Binder holes
    y = 2.0
    for _ in range(3):
        hole = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.62), Inches(y), Inches(0.12), Inches(0.12))
        hole.fill.solid()
        hole.fill.fore_color.rgb = RGBColor(0xC9, 0xCD, 0xD2)
        hole.line.fill.background()
        y += 1.0

    # Right tabs
    tabs = [COLOR_TAB_1, COLOR_TAB_2, COLOR_TAB_3, COLOR_TAB_4, COLOR_TAB_5]
    y = 1.1
    for c in tabs:
        tab = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(12.68),
            Inches(y),
            Inches(0.45),
            Inches(0.78),
        )
        tab.fill.solid()
        tab.fill.fore_color.rgb = c
        tab.line.fill.background()
        y += 0.95


def add_meta(slide, left: str, right: str) -> None:
    box = slide.shapes.add_textbox(Inches(1.0), Inches(0.75), Inches(11.1), Inches(0.35))
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = left
    run.font.name = "DM Sans"
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_SOFT

    p.alignment = PP_ALIGN.LEFT

    box2 = slide.shapes.add_textbox(Inches(8.6), Inches(0.75), Inches(3.5), Inches(0.35))
    tf2 = box2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = right
    run2.font.name = "DM Sans"
    run2.font.size = Pt(11)
    run2.font.color.rgb = COLOR_SOFT


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(10.9), Inches(1.05))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Bodoni Moda"
    run.font.bold = True
    run.font.size = Pt(38)
    run.font.color.rgb = COLOR_INK


def add_subtitle(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(10.7), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "DM Sans"
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_SOFT


def add_bullets(slide, title: str, bullets: list[str]) -> None:
    add_title(slide, title)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.0), Inches(4.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "DM Sans"
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_SOFT
        p.space_after = Pt(12)
        p.bullet = True


def add_cards(slide, title: str, cards: list[tuple[str, str]]) -> None:
    add_title(slide, title)
    x_positions = [1.0, 4.25, 7.5]
    y_positions = [2.3, 4.35]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            if idx >= len(cards):
                return
            h, body = cards[idx]
            idx += 1
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(3.05),
                Inches(1.85),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card.line.color.rgb = RGBColor(0xDF, 0xD9, 0xCE)
            card.line.width = Pt(1)

            t = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.12), Inches(2.7), Inches(0.45))
            tf_t = t.text_frame
            tf_t.clear()
            p_t = tf_t.paragraphs[0]
            p_t.text = h
            p_t.font.name = "Bodoni Moda"
            p_t.font.bold = True
            p_t.font.size = Pt(16)
            p_t.font.color.rgb = COLOR_INK

            b = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.64), Inches(2.7), Inches(1.08))
            tf_b = b.text_frame
            tf_b.word_wrap = True
            tf_b.clear()
            p_b = tf_b.paragraphs[0]
            p_b.text = body
            p_b.font.name = "DM Sans"
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = COLOR_SOFT


def add_process(slide, title: str, steps: list[tuple[str, str]]) -> None:
    add_title(slide, title)
    x_positions = [1.0, 4.25, 7.5]
    y_positions = [2.3, 4.35]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            if idx >= len(steps):
                return
            label, body = steps[idx]
            idx += 1
            box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(3.05),
                Inches(1.85),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF1, 0xF7, 0xFB)
            box.line.color.rgb = COLOR_ACCENT
            box.line.width = Pt(1)

            t = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.12), Inches(2.75), Inches(0.42))
            tf_t = t.text_frame
            tf_t.clear()
            p_t = tf_t.paragraphs[0]
            p_t.text = label
            p_t.font.name = "DM Sans"
            p_t.font.bold = True
            p_t.font.size = Pt(12)
            p_t.font.color.rgb = COLOR_INK

            b = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.56), Inches(2.75), Inches(1.15))
            tf_b = b.text_frame
            tf_b.word_wrap = True
            tf_b.clear()
            p_b = tf_b.paragraphs[0]
            p_b.text = body
            p_b.font.name = "DM Sans"
            p_b.font.size = Pt(10)
            p_b.font.color.rgb = COLOR_SOFT


def add_table_like(slide, title: str) -> None:
    add_title(slide, title)

    cols = ["Item", "AWS Instance", "Azure SKU", "Region", "AWS PayGo", "Azure PayGo"]
    rows = [
        ["row-1", "c6i.xlarge", "Standard_F4s_v2", "Sydney / australiaeast", "0.222", "0.222"],
        ["row-2", "c6i.xlarge", "Standard_F4s_v2", "Tokyo / japaneast", "0.214", "0.214"],
        ["row-5", "t2.small", "Standard_B2pts_v2", "Oregon / westus2", "0.023", "0.0084"],
        ["row-6", "r7a.2xlarge", "Standard_E8s_v5", "Oregon / westus2", "0.6086", "0.504"],
    ]

    table_shape = slide.shapes.add_table(5, 6, Inches(1.0), Inches(2.2), Inches(11.1), Inches(2.5))
    table = table_shape.table

    widths = [1.0, 2.0, 2.3, 2.4, 1.6, 1.8]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for c, col in enumerate(cols):
        cell = table.cell(0, c)
        cell.text = col
        p = cell.text_frame.paragraphs[0]
        p.font.name = "DM Sans"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_INK

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.name = "DM Sans"
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_SOFT

    note = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(10.9), Inches(1.0))
    tf = note.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "样例保留 review_reason 与 evidence.source_ref，支持审计回溯与人工复核。"
    p.font.name = "DM Sans"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_SOFT


def add_code_block(slide, title: str, code: str) -> None:
    add_title(slide, title)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(2.2),
        Inches(11.1),
        Inches(2.9),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x13, 0x17, 0x1D)
    box.line.fill.background()

    txt = slide.shapes.add_textbox(Inches(1.25), Inches(2.45), Inches(10.6), Inches(2.35))
    tf = txt.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xD6, 0xDE, 0xE8)

    note = slide.shapes.add_textbox(Inches(1.0), Inches(5.35), Inches(11.1), Inches(0.9))
    tf2 = note.text_frame
    tf2.word_wrap = True
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = "支持按 skill 粒度执行映射与定价；maintenance 脚本用于刷新 AWS 全区域报价文件。"
    p2.font.name = "DM Sans"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_SOFT


def add_quote_slide(slide, quote: str, author: str) -> None:
    add_meta(slide, "Closing", "Decision-Ready Quoting")
    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.8), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"\"{quote}\""
    p.font.name = "Bodoni Moda"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_INK

    by = slide.shapes.add_textbox(Inches(1.2), Inches(4.65), Inches(6), Inches(0.6))
    tf2 = by.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = author
    p2.font.name = "DM Sans"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_SOFT


def create_intro_deck(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1)
    add_sheet(s1)
    add_meta(s1, "Azure Quote Agent", "Project Introduction Deck")
    add_title(s1, "把混乱输入变成可决策的迁移报价")
    add_subtitle(s1, "面向云迁移场景，从 PDF/Excel/CSV 中抽取 VM 需求，映射 Azure 机型，拉取双云定价，并输出结构化报价载体。")

    # 2 Problem
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2)
    add_sheet(s2)
    add_meta(s2, "Problem", "Why Build It")
    add_bullets(
        s2,
        "手工迁移报价的三大痛点",
        [
            "输入来源异构：账单 PDF、整理 Excel、半结构化 CSV 混杂。",
            "映射逻辑复杂：规格、区域、SAP 负载、候选回退策略同时耦合。",
            "价格口径差异：AWS/Azure 定价来源与 RI 口径不同。",
            "交付格式不一致：经常要求 JSON + Excel + 可追溯证据。",
            "结果可信度难说明：review 项目难以快速定位。",
        ],
    )

    # 3 Value
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3)
    add_sheet(s3)
    add_meta(s3, "Outcome", "What You Get")
    add_cards(
        s3,
        "一个端到端、可追溯的报价流水线",
        [
            ("3 类输入", "PDF / Excel / CSV 统一归一化"),
            ("2 云定价", "AWS + Azure PayGo/RI 对比"),
            ("7 个 MCP 工具", "可被 Agent 直接编排调用"),
            ("3 类输出", "CSV / quote_payload / Excel"),
            ("CLI + MCP", "既可脚本运行，也可服务化接入"),
            ("Review 机制", "不确定项自动聚合与标注"),
        ],
    )

    # 4 Pipeline
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4)
    add_sheet(s4)
    add_meta(s4, "Architecture", "Data Pipeline")
    add_process(
        s4,
        "核心处理链路（6 步）",
        [
            ("01 输入抽取", "Excel 标准化 + PDF/DI 解析"),
            ("02 区域归一", "自然语言地域映射为 canonical region"),
            ("03 规格识别", "解析 vCPU/内存/架构与家族特征"),
            ("04 SAP 推断", "基于 system/env/workload_type 规则判定"),
            ("05 报价比价", "拉取 AWS/Azure PayGo、1Y/3Y RI"),
            ("06 交付封装", "生成 quote_payload 与报价 Excel"),
        ],
    )

    # 5 Modules
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5)
    add_sheet(s5)
    add_meta(s5, "Codebase", "scripts/ Modules")
    add_cards(
        s5,
        "关键模块地图",
        [
            ("extract_excel_inputs.py", "多列别名与 profile 抽取"),
            ("pdf_extraction_core.py", "Azure DI 解析与行级过滤"),
            ("region_mapping_core.py", "区域映射与置信度输出"),
            ("sap_inference.py", "SAP 负载规则推断"),
            ("build_vm_quote_payload.py", "汇总结构化 payload"),
            ("write_quote_excel.py", "写入可交付报价工作簿"),
        ],
    )

    # 6 MCP
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6)
    add_sheet(s6)
    add_meta(s6, "MCP Server", "Tool Surface")
    add_cards(
        s6,
        "Agent 可调用能力（mcp_server.py）",
        [
            ("map_region_single", "单条地域映射 + 置信度"),
            ("map_region_batch", "批量映射与命中率统计"),
            ("map_region_file", "文件级映射任务"),
            ("extract_pdf_inputs", "单 PDF 解析并输出 CSV"),
            ("extract_pdf_inputs_batch", "多 PDF 批处理与汇总"),
            ("validate_di_connection", "凭证连通性与端到端预检"),
        ],
    )

    # 7 Sample table
    s7 = prs.slides.add_slide(blank)
    set_slide_bg(s7)
    add_sheet(s7)
    add_meta(s7, "Evidence", "output/sample_input_4_pdf")
    add_table_like(s7, "样例结果展示：定价与映射落地")

    # 8 Payload model
    s8 = prs.slides.add_slide(blank)
    set_slide_bg(s8)
    add_sheet(s8)
    add_meta(s8, "Data Contract", "quote_payload.json")
    add_bullets(
        s8,
        "统一交付结构：4 个顶层对象",
        [
            "summary：客户项目、币种、竞品云、价格日期与来源说明。",
            "line_items：规格、区域、候选 SKU、双云价格、review 字段。",
            "assumptions：自动归纳复核假设，标记影响范围。",
            "evidence：source_url、fetched_at、source_ref 可追溯。",
            "可直接输入 write_quote_excel.py 生成四工作表交付件。",
        ],
    )

    # 9 CLI
    s9 = prs.slides.add_slide(blank)
    set_slide_bg(s9)
    add_sheet(s9)
    add_meta(s9, "Usage", "CLI + MCP")
    add_code_block(
        s9,
        "典型调用路径（示意）",
        "python scripts/extract_excel_inputs.py --input-excel input/sample.xlsx --output output/extracted_inputs.csv\n"
        "python scripts/build_vm_quote_payload.py --input-csv output/vm_pricing_results.csv --output-json output/quote_payload.json\n"
        "python scripts/write_quote_excel.py --input-json output/quote_payload.json --output-xlsx output/quote.xlsx\n"
        "python scripts/mcp_server.py",
    )

    # 10 Quality
    s10 = prs.slides.add_slide(blank)
    set_slide_bg(s10)
    add_sheet(s10)
    add_meta(s10, "Reliability", "Review Mechanism")
    add_bullets(
        s10,
        "质量保障：自动标记 + 证据追踪",
        [
            "review_flag/review_reason 将不确定项显式分离。",
            "build_dynamic_review_assumptions 自动聚合同类问题。",
            "异常时可保留可用结果，并将风险落入 assumptions。",
            "MCP 层限制文件路径在 input/output/data 白名单内。",
            "validate_di_connection 降低批处理任务失败率。",
        ],
    )

    # 11 Roadmap
    s11 = prs.slides.add_slide(blank)
    set_slide_bg(s11)
    add_sheet(s11)
    add_meta(s11, "Next", "Potential Enhancements")
    add_cards(
        s11,
        "下一步可扩展方向",
        [
            ("自动化场景库", "沉淀行业模板，一键生成报价参数"),
            ("观测性面板", "命中率、fallback 率、错误分布可视化"),
            ("并发优化", "缓存零售价格与 offer 文件，降低时延"),
            ("策略配置化", "SKU 权重与门禁规则从代码抽离"),
            ("审计增强", "payload 签名与版本锁，便于合规留档"),
            ("多语言输出", "同一 payload 生成中英文材料"),
        ],
    )

    # 12 Closing
    s12 = prs.slides.add_slide(blank)
    set_slide_bg(s12)
    add_sheet(s12)
    add_quote_slide(s12, "让迁移报价从脚本拼接，升级为可复用、可验证、可交付的工程化产品。", "Azure Quote Agent")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    project_root = Path(__file__).resolve().parents[1]
    output = project_root / "output" / "azure-quote-agent-introduction-deck-editable.pptx"
    create_intro_deck(output)
    print(f"Editable PPTX generated: {output.as_posix()}")


if __name__ == "__main__":
    main()
