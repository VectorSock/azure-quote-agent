from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors (Ocean Gradient)
COLOR_PRIMARY = RGBColor(0x06, 0x5A, 0x82)   # Deep Blue
COLOR_SECONDARY = RGBColor(0x1C, 0x72, 0x93) # Teal
COLOR_ACCENT = RGBColor(0x21, 0x29, 0x5C)    # Midnight
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT_DARK = RGBColor(0x36, 0x45, 0x4F) # Charcoal

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Customize background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Azure Quote Agent"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Automated AWS to Azure Migration Pricing & Wrapping"
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    subtitle.text_frame.paragraphs[0].font.name = "Arial"

    # --- Slide 2: Overview ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Overview"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A comprehensive automation tool designed to streamline the migration quote process from AWS to Azure."
    
    p = tf.add_paragraph()
    p.text = "Key Capabilities:"
    p.level = 0
    p.font.bold = True
    
    bullets = [
        "Parses AWS configurations from PDF invoices and CSV files.",
        "Maps AWS instances to Azure VM equivalents automatically.",
        "Detects SAP specialized workloads.",
        "Generates structured quote payloads via MCP.",
    ]
    
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1

    # --- Slide 3: System Architecture (Diagrammatic) ---
    slide_layout = prs.slide_layouts[5] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title = title_shape.text_frame
    title.text = "System Architecture"
    title.paragraphs[0].font.size = Pt(40)
    title.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title.paragraphs[0].font.bold = True

    # Draw boxes for components
    def add_box(text, x, y, w, h, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.text_frame.text = text
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return shape

    # Input Layer
    add_box("Input Sources\n(PDF / CSV)", 1, 2.5, 2, 1.5, COLOR_SECONDARY)
    
    # Processing Layer
    add_box("Core Processing\n(Python Scripts)", 4, 2, 3, 2.5, COLOR_PRIMARY)
    
    # Internal Logic
    slide.shapes.add_textbox(Inches(4.2), Inches(2.2), Inches(2.5), Inches(0.5)).text_frame.text = "• PDF Extraction (Azure DI)"
    slide.shapes.add_textbox(Inches(4.2), Inches(2.7), Inches(2.5), Inches(0.5)).text_frame.text = "• Region Mapping"
    slide.shapes.add_textbox(Inches(4.2), Inches(3.2), Inches(2.5), Inches(0.5)).text_frame.text = "• SAP Inference"

    # MCP Layer
    add_box("MCP Server\n(Interface)", 4, 5, 3, 1, COLOR_ACCENT)

    # Output Layer
    add_box("Quote Payload\n(JSON / Excel)", 8, 2.5, 2, 1.5, COLOR_SECONDARY)

    # Arrows (simplified lines)
    # Using shapes instead of connectors for visual arrows
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(3.5), Inches(0.8), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_TEXT_DARK
    
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.1), Inches(3.5), Inches(0.8), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_TEXT_DARK


    # --- Slide 4: Key Components ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Components"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    components = [
        ("PDF Extraction", "Uses Azure Document Intelligence to convert unstructured invoice data into structured records."),
        ("Region Mapping", "Intelligent resolution of AWS regions (e.g., 'us-west-2') to Azure equivalents (e.g., 'West US 2')."),
        ("SAP Inference", "Analyzes workload patterns (HANA, NetWeaver) to recommend certified Azure instances."),
        ("MCP Server", "Exposes all functionality through the Model Context Protocol for Agent integration."),
    ]
    
    for head, desc in components:
        p = tf.add_paragraph()
        p.text = head
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.level = 1

    # --- Slide 5: Usage & Integration ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Usage & Integration"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The tool is designed to be used both standalone and as an agent tool."
    
    p = tf.add_paragraph()
    p.text = "Command Line:"
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "python scripts/build_vm_quote_payload.py --input <file>"
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "MCP Integration:"
    p.font.bold = True
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Running 'mcp_server.py' exposes 'parse_pdf' and 'calculate_quote' tools to AI agents."
    p.level = 1

    # --- Slide 6: Conclusion ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_ACCENT

    title = slide.shapes.title
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Azure Quote Agent accelerates the migration estimation process with accuracy and intelligence."
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY


    output_path = "Azure_Quote_Agent_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    try:
        create_presentation()
        print("Presentation created successfully!")
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"FAILED: {e}")
