from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
OUTPUT_PPT = OUT_DIR / "Azure-Quote-Agent-Project-Overview.pptx"


def safe_read_csv(path: Path) -> pd.DataFrame:
    if not path.exists():
        return pd.DataFrame()
    return pd.read_csv(path)


def build_stats() -> dict:
    extracted = safe_read_csv(OUT_DIR / "extracted_inputs.csv")
    regions = safe_read_csv(OUT_DIR / "extracted_with_regions.csv")
    mapping = safe_read_csv(OUT_DIR / "azure_instance_mapping_sample4.csv")
    pricing = safe_read_csv(OUT_DIR / "vm_pricing_results_sample4.csv")

    total_inputs = len(extracted)
    mapped_region_rows = 0
    unique_azure_regions = 0
    if not regions.empty and "mapped_azure_region" in regions.columns:
        mapped_region_rows = int(regions["mapped_azure_region"].notna().sum())
        unique_azure_regions = int(regions["mapped_azure_region"].dropna().nunique())

    mapped_sku_rows = 0
    fallback_rows = 0
    review_rows = 0
    if not mapping.empty:
        if "azure_sku" in mapping.columns:
            mapped_sku_rows = int(mapping["azure_sku"].notna().sum())
        if "fallback_sku" in mapping.columns:
            fallback_rows = int(mapping["fallback_sku"].notna().sum())
        if "review_flag" in mapping.columns:
            review_rows = int(pd.to_numeric(mapping["review_flag"], errors="coerce").fillna(0).astype(bool).sum())

    pricing_ok_rows = 0
    azure_not_found_rows = 0
    aws_priced_rows = 0
    avg_azure_paygo = 0.0
    if not pricing.empty:
        if "pricing_status" in pricing.columns:
            pricing_ok_rows = int((pricing["pricing_status"].astype(str) == "ok").sum())
        if "azure_status" in pricing.columns:
            azure_not_found_rows = int((pricing["azure_status"].astype(str) == "not_found").sum())
        if "aws_status" in pricing.columns:
            aws_priced_rows = int((pricing["aws_status"].astype(str) == "ok").sum())
        if "Azure_paygo" in pricing.columns:
            paygo_series = pd.to_numeric(pricing["Azure_paygo"], errors="coerce").dropna()
            avg_azure_paygo = float(paygo_series.mean()) if not paygo_series.empty else 0.0

    return {
        "total_inputs": total_inputs,
        "mapped_region_rows": mapped_region_rows,
        "unique_azure_regions": unique_azure_regions,
        "mapped_sku_rows": mapped_sku_rows,
        "fallback_rows": fallback_rows,
        "review_rows": review_rows,
        "pricing_ok_rows": pricing_ok_rows,
        "azure_not_found_rows": azure_not_found_rows,
        "aws_priced_rows": aws_priced_rows,
        "avg_azure_paygo": avg_azure_paygo,
    }


def add_title(slide, title: str, subtitle: str = ""):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 42, 74)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.2), Inches(11.8), Inches(0.6))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(55, 76, 98)


def add_bullets(slide, items: list[str], x=0.9, y=1.9, w=11.2, h=4.8, size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(35, 35, 35)
        p.space_after = Pt(10)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.RIGHT


def make_presentation(stats: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s1, "Azure Quote Agent", "Project Overview")
    add_bullets(
        s1,
        [
            "Automates cross-cloud VM quotation from raw request sheets to final quote workbook",
            "Pipeline: input extraction -> normalization -> SKU mapping -> pricing -> deliverable export",
            f"Snapshot: {stats['total_inputs']} VM rows processed in current sample outputs",
        ],
        size=24,
        y=2.0,
        h=3.6,
    )
    add_footer(s1, f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s2, "Business Problem")
    add_bullets(
        s2,
        [
            "Source requests are messy: Excel/PDF with mixed providers, fuzzy regions, and missing fields",
            "Manual SKU matching and price checks are slow and inconsistent across teams",
            "Need a repeatable, evidence-backed quote artifact for customer-facing proposals",
        ],
    )

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "End-to-End Architecture")
    add_bullets(
        s3,
        [
            "1. Input extraction: parse Excel/PDF into canonical columns",
            "2. Region normalization: map city/region text to AWS/Azure/GCP region IDs",
            "3. Workload inference: derive SAP and workload hints from context",
            "4. Instance mapping: map VM config to Azure SKU with support gates and ranking",
            "5. Pricing: fetch Azure/AWS retail + reservation pricing",
            "6. Quote output: build payload and write Excel (Summary/LineItems/Assumptions/Evidence)",
        ],
        size=19,
        y=1.7,
        h=4.8,
    )

    # Slide 4
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "Core Modules")
    add_bullets(
        s4,
        [
            "input-excel-extraction / input-pdf-extraction: structural data capture",
            "global-region-mapping: fuzzy location to cloud regions",
            "vm-aws-instance-to-config: AWS instance parsing and indicator inference",
            "shared/sap_inference.py: SAP workload classification",
            "vm-config-to-azure-instance: Azure SKU candidate generation and ranking",
            "vm-pricing-retail-api: Azure + AWS price retrieval with fallback",
            "global-quote-writer: quote payload + final XLSX generation",
        ],
        size=18,
        y=1.6,
        h=5.2,
    )

    # Slide 5
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s5, "Data Contract and Governance")
    add_bullets(
        s5,
        [
            "Canonical schema defined in .github/skills/references/column-schema.md",
            "Unified fields include: resource specs, mapped regions, inferred workload, mapped SKU, pricing",
            "Each stage appends fields while preserving upstream provenance for traceability",
            "Status/reason columns make review and exception handling explicit",
        ],
        size=20,
    )

    # Slide 6
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s6, "Sample Run Metrics")
    metric_lines = [
        f"Input rows: {stats['total_inputs']}",
        f"Rows with mapped Azure region: {stats['mapped_region_rows']}",
        f"Unique mapped Azure regions: {stats['unique_azure_regions']}",
        f"Rows with Azure SKU mapping: {stats['mapped_sku_rows']}",
        f"Pricing status=ok rows: {stats['pricing_ok_rows']}",
        f"Azure pricing not_found rows: {stats['azure_not_found_rows']}",
        f"AWS priced rows: {stats['aws_priced_rows']}",
        f"Average Azure PAYG hourly price: ${stats['avg_azure_paygo']:.3f}",
    ]
    add_bullets(s6, metric_lines, size=21, y=1.8, h=4.8)

    # Slide 7
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s7, "Mapping and Ranking Logic")
    add_bullets(
        s7,
        [
            "Support gate filters out invalid candidates (region, capability, policy constraints)",
            "Ranking score blends fit, performance hints, and cost signal",
            "Fallback strategy keeps alternatives when primary is unavailable",
            "SAP-aware path can route to certified SKUs when workload indicates SAP",
        ],
        size=20,
    )

    # Slide 8
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s8, "Pricing Strategy")
    add_bullets(
        s8,
        [
            "Azure prices from Retail Prices API with generation-aware SKU selection",
            "AWS prices from Pricing API and offer-file fallback when needed",
            "Outputs include PAYG, 1Y RI, and 3Y RI hourly prices",
            "Per-row JSON evidence captured for auditability and troubleshooting",
        ],
        size=20,
    )

    # Slide 9
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s9, "Deliverables Produced")
    add_bullets(
        s9,
        [
            "Intermediate CSVs: extracted_inputs, extracted_with_regions, mapping, pricing",
            "Quote payload JSON for downstream integration",
            "Final Excel quote workbook with 4 tabs:",
            "Summary, LineItems, Assumptions, Evidence",
        ],
        size=21,
    )

    # Slide 10
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s10, "Strengths and Next Steps")
    add_bullets(
        s10,
        [
            "Strengths: modular pipeline, transparent assumptions, cross-cloud coverage",
            "Add stronger validation for missing/ambiguous regions",
            "Introduce automated regression tests per skill module",
            "Add cost optimization recommendations after baseline mapping",
            "Package full flow into a single orchestrated CLI command",
        ],
        size=19,
    )

    return prs


def main() -> None:
    stats = build_stats()
    prs = make_presentation(stats)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPT)
    print(f"Saved: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
